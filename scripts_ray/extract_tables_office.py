#!/usr/bin/env python3
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import os
from pathlib import Path

class OfficeTableExtractor:
    """Office 文档表格提取器"""
    
    def __init__(self, output_format='csv'):
        """
        output_format: 'csv', 'markdown', 'excel', 'html'
        """
        self.output_format = output_format
    
    def extract_all(self, directory=None):
        """提取目录下所有 Office 文档的表格"""
        if directory is None:
            directory = Path.cwd()
        else:
            directory = Path(directory)
        
        # 处理 DOCX
        for docx_file in directory.glob("*.docx"):
            if not docx_file.name.startswith('~$'):
                self.extract_from_docx(docx_file)
        
        # 处理 PPTX
        for pptx_file in directory.glob("*.pptx"):
            if not pptx_file.name.startswith('~$'):
                self.extract_from_pptx(pptx_file)
    
    def extract_from_docx(self, docx_path):
        """从 DOCX 提取表格"""
        print(f"\n处理 DOCX: {docx_path.name}")
        doc = Document(str(docx_path))
        output_dir = Path(f"{docx_path.stem}_tables")
        output_dir.mkdir(exist_ok=True)
        
        for i, table in enumerate(doc.tables):
            data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                data.append(row_data)
            
            if data:
                self._save_table(data, output_dir, f"table_{i+1}", i+1)
    
    def extract_from_pptx(self, pptx_path):
        """从 PPTX 提取表格"""
        print(f"\n处理 PPTX: {pptx_path.name}")
        prs = Presentation(str(pptx_path))
        output_dir = Path(f"{pptx_path.stem}_tables")
        output_dir.mkdir(exist_ok=True)
        
        table_count = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    table_count += 1
                    data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        data.append(row_data)
                    
                    if data:
                        filename = f"slide{slide_num}_table{table_count}"
                        self._save_table(data, output_dir, filename, table_count)
    
    def _save_table(self, data, output_dir, filename, table_num):
        """根据指定格式保存表格"""
        if self.output_format == 'csv':
            output_file = output_dir / f"{filename}.csv"
            df = pd.DataFrame(data[1:], columns=data[0] if data else [])
            df.to_csv(output_file, index=False, encoding='utf-8')
            
        elif self.output_format == 'markdown':
            output_file = output_dir / f"{filename}.md"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(self._to_markdown(data))
                
        elif self.output_format == 'excel':
            output_file = output_dir / f"{filename}.xlsx"
            df = pd.DataFrame(data[1:], columns=data[0] if data else [])
            df.to_excel(output_file, index=False)
            
        elif self.output_format == 'html':
            output_file = output_dir / f"{filename}.html"
            df = pd.DataFrame(data[1:], columns=data[0] if data else [])
            df.to_html(output_file, index=False)
        
        print(f"  已保存表格 {table_num}: {output_file}")
    
    def _to_markdown(self, data):
        """转换为 Markdown 格式"""
        if not data:
            return ""
        
        md = []
        # 表头
        md.append("| " + " | ".join(data[0]) + " |")
        md.append("|" + "|".join([" --- " for _ in data[0]]) + "|")
        
        # 数据行
        for row in data[1:]:
            md.append("| " + " | ".join(row) + " |")
        
        return "\n".join(md)

# 使用示例
if __name__ == "__main__":
    # 提取为 CSV
    extractor = OfficeTableExtractor(output_format='csv')
    extractor.extract_all()
    
    # 提取为 Markdown
    # extractor = OfficeTableExtractor(output_format='markdown')
    # extractor.extract_all()

