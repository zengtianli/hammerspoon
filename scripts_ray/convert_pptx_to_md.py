#!/usr/bin/env python3
"""
PPTX to Markdown Converter (Fixed Version)
修复版：兼容不同版本的 python-pptx
"""

import os
import sys
import argparse
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import warnings

# 忽略一些无关紧要的警告
warnings.filterwarnings("ignore")

# 检查 MSO_SHAPE_TYPE 中的可用属性
SHAPE_TYPES = {
    'PICTURE': getattr(MSO_SHAPE_TYPE, 'PICTURE', 13),
    'PLACEHOLDER': getattr(MSO_SHAPE_TYPE, 'PLACEHOLDER', 14),
    'TEXT_BOX': getattr(MSO_SHAPE_TYPE, 'TEXT_BOX', 17),
    'TABLE': getattr(MSO_SHAPE_TYPE, 'TABLE', 19),
    'CHART': getattr(MSO_SHAPE_TYPE, 'CHART', 3),
    'GROUP': getattr(MSO_SHAPE_TYPE, 'GROUP', 6),
}

class PPTXToMarkdownConverter:
    def __init__(self, extract_images=True, verbose=False):
        self.extract_images = extract_images
        self.verbose = verbose
        self.image_counter = 0
        self.table_counter = 0
        
    def log(self, message):
        """打印详细信息"""
        if self.verbose:
            print(f"  {message}")
    
    def clean_text(self, text):
        """清理文本，移除多余的空白字符"""
        if not text:
            return ""
        # 移除多余的空白行
        lines = text.split('\n')
        cleaned_lines = [line.rstrip() for line in lines]
        # 移除开头和结尾的空行
        while cleaned_lines and not cleaned_lines[0]:
            cleaned_lines.pop(0)
        while cleaned_lines and not cleaned_lines[-1]:
            cleaned_lines.pop()
        return '\n'.join(cleaned_lines)
    
    def extract_text_from_shape(self, shape):
        """从形状中提取格式化文本"""
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            return ""
        
        formatted_text = []
        text_frame = shape.text_frame
        
        for paragraph in text_frame.paragraphs:
            para_text = ""
            
            # 检查段落级别（用于缩进）
            level = paragraph.level if hasattr(paragraph, 'level') else 0
            indent = "  " * level
            
            # 提取段落中的文本
            for run in paragraph.runs:
                text = run.text
                if not text:
                    continue
                
                # 检查格式
                try:
                    if run.font.bold:
                        text = f"**{text}**"
                    if run.font.italic:
                        text = f"*{text}*"
                    if run.font.underline:
                        text = f"<u>{text}</u>"
                except:
                    # 如果无法获取字体属性，使用原始文本
                    pass
                
                para_text += text
            
            if para_text:
                # 检查是否是列表项
                bullet = paragraph.text.strip()
                if bullet and bullet[0] in '•·●○■□▪▫◆◇★☆▸▹►▻':
                    formatted_text.append(f"{indent}- {para_text}")
                elif re.match(r'^\d+[\.\)]\s', paragraph.text):
                    # 编号列表
                    formatted_text.append(f"{indent}{para_text}")
                else:
                    formatted_text.append(f"{indent}{para_text}")
        
        return '\n'.join(formatted_text)
    
    def extract_table_from_shape(self, shape):
        """从表格形状中提取内容并转换为 Markdown 表格"""
        if not hasattr(shape, 'has_table') or not shape.has_table:
            return ""
        
        self.table_counter += 1
        self.log(f"找到表格 #{self.table_counter}")
        
        try:
            table = shape.table
            md_table = []
            
            # 获取最大列数
            max_cols = max(len(row.cells) for row in table.rows) if table.rows else 0
            
            # 提取表格内容
            for row_idx, row in enumerate(table.rows):
                row_data = []
                
                for col_idx in range(max_cols):
                    if col_idx < len(row.cells):
                        cell = row.cells[col_idx]
                        # 清理单元格文本
                        cell_text = cell.text.strip().replace('\n', ' ')
                        # 转义管道符
                        cell_text = cell_text.replace('|', '\\|')
                        row_data.append(cell_text)
                    else:
                        row_data.append("")
                
                md_table.append('| ' + ' | '.join(row_data) + ' |')
                
                # 在第一行后添加分隔符
                if row_idx == 0:
                    separator = '|'
                    for _ in range(max_cols):
                        separator += ' --- |'
                    md_table.append(separator)
            
            return '\n'.join(md_table)
        except Exception as e:
            self.log(f"提取表格失败: {e}")
            return "[表格]"
    
    def save_image(self, shape, output_dir, slide_num):
        """保存图片并返回相对路径"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # 确定图片格式
            ext = image.ext or 'png'
            
            # 创建图片目录
            img_dir = Path(output_dir) / 'images'
            img_dir.mkdir(exist_ok=True, parents=True)
            
            # 生成唯一的图片文件名
            self.image_counter += 1
            img_filename = f'slide{slide_num:03d}_img{self.image_counter:03d}.{ext}'
            img_path = img_dir / img_filename
            
            # 保存图片
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            
            self.log(f"保存图片: {img_filename}")
            
            # 返回相对路径
            return f'images/{img_filename}'
            
        except Exception as e:
            self.log(f"⚠️  保存图片失败: {e}")
            return None
    
    def get_shape_type_name(self, shape):
        """获取形状类型名称"""
        try:
            shape_type = shape.shape_type
            
            # 检查常见类型
            if shape_type == SHAPE_TYPES['PICTURE']:
                return 'PICTURE'
            elif shape_type == SHAPE_TYPES['TABLE']:
                return 'TABLE'
            elif shape_type == SHAPE_TYPES['CHART']:
                return 'CHART'
            elif shape_type == SHAPE_TYPES['TEXT_BOX']:
                return 'TEXT_BOX'
            elif shape_type == SHAPE_TYPES['PLACEHOLDER']:
                return 'PLACEHOLDER'
            elif shape_type == SHAPE_TYPES['GROUP']:
                return 'GROUP'
            else:
                return f'OTHER_{shape_type}'
        except:
            return 'UNKNOWN'
    
    def convert_slide_to_markdown(self, slide, slide_num, output_dir):
        """将单个幻灯片转换为 Markdown"""
        md_content = []
        
        # 添加幻灯片分隔符
        md_content.append("\n---")
        
        # 添加幻灯片编号
        md_content.append(f"\n## 幻灯片 {slide_num}")
        
        # 获取幻灯片备注
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    md_content.append(f"\n> **备注**: {notes.strip()}")
        except:
            pass
        
        md_content.append("")
        
        # 提取标题
        title_text = ""
        try:
            if slide.shapes.title:
                title_text = self.clean_text(slide.shapes.title.text)
                if title_text:
                    md_content.append(f"### {title_text}")
                    md_content.append("")
        except:
            pass
        
        # 按位置排序形状（从上到下，从左到右）
        shapes_with_pos = []
        for shape in slide.shapes:
            try:
                if shape == slide.shapes.title:
                    continue
                
                top = shape.top.cm if hasattr(shape, 'top') and hasattr(shape.top, 'cm') else 0
                left = shape.left.cm if hasattr(shape, 'left') and hasattr(shape.left, 'cm') else 0
                shapes_with_pos.append((top, left, shape))
            except:
                # 如果无法获取位置，仍然添加形状
                shapes_with_pos.append((0, 0, shape))
        
        # 排序
        shapes_with_pos.sort(key=lambda x: (x[0], x[1]))
        
        # 处理每个形状
        for _, _, shape in shapes_with_pos:
            content_added = False
            
            try:
                # 获取形状类型
                shape_type_name = self.get_shape_type_name(shape)
                self.log(f"处理形状: {shape_type_name}")
                
                # 处理文本框
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = self.extract_text_from_shape(shape)
                    if text.strip():
                        md_content.append(text)
                        content_added = True
                
                # 处理表格
                elif hasattr(shape, 'has_table') and shape.has_table:
                    table_md = self.extract_table_from_shape(shape)
                    if table_md:
                        md_content.append("")
                        md_content.append(table_md)
                        content_added = True
                
                # 处理图片
                elif shape_type_name == 'PICTURE':
                    if self.extract_images:
                        img_path = self.save_image(shape, output_dir, slide_num)
                        if img_path:
                            alt_text = f"Slide {slide_num} - Image {self.image_counter}"
                            md_content.append(f"\n![{alt_text}]({img_path})")
                        else:
                            md_content.append(f"\n[图片无法提取]")
                    else:
                        md_content.append(f"\n[图片：幻灯片 {slide_num}]")
                    content_added = True
                
                # 处理图表
                elif hasattr(shape, 'has_chart') and shape.has_chart:
                    chart_title = "未命名图表"
                    try:
                        if shape.chart.has_title:
                            chart_title = shape.chart.chart_title.text_frame.text
                    except:
                        pass
                    md_content.append(f"\n[图表：{chart_title}]")
                    content_added = True
                
                # 处理组合形状
                elif shape_type_name == 'GROUP':
                    md_content.append(f"\n[组合形状]")
                    # 可以递归处理组内的形状
                    content_added = True
                
                # 处理其他类型
                elif shape_type_name not in ['TEXT_BOX', 'PLACEHOLDER']:
                    md_content.append(f"\n[{shape_type_name}]")
                    content_added = True
                
            except Exception as e:
                self.log(f"处理形状时出错: {e}")
                continue
            
            if content_added:
                md_content.append("")
        
        return '\n'.join(md_content)
    
    def convert_presentation(self, pptx_path, output_path=None):
        """转换整个演示文稿"""
        try:
            print(f"\n📂 正在处理: {pptx_path}")
            
            # 加载演示文稿
            prs = Presentation(pptx_path)
            
            # 创建与PPTX文件同名的输出目录
            pptx_name = Path(pptx_path).stem
            output_dir = Path(pptx_path).parent / pptx_name
            output_dir.mkdir(exist_ok=True)
            
            # 确定输出路径
            if output_path is None:
                output_path = output_dir / f"{pptx_name}.md"
            
            # 重置计数器
            self.image_counter = 0
            self.table_counter = 0
            
            # 构建 Markdown 内容
            md_content = []
            
            # 添加文档标题
            md_content.append(f"# {Path(pptx_path).stem}")
            md_content.append("")
            
            # 添加元信息
            md_content.append("## 文档信息")
            md_content.append("")
            md_content.append(f"- **源文件**: `{Path(pptx_path).name}`")
            md_content.append(f"- **幻灯片数量**: {len(prs.slides)}")
            
            # 尝试获取演示文稿属性
            try:
                core_props = prs.core_properties
                if hasattr(core_props, 'author') and core_props.author:
                    md_content.append(f"- **作者**: {core_props.author}")
                if hasattr(core_props, 'title') and core_props.title:
                    md_content.append(f"- **标题**: {core_props.title}")
                if hasattr(core_props, 'subject') and core_props.subject:
                    md_content.append(f"- **主题**: {core_props.subject}")
                if hasattr(core_props, 'created') and core_props.created:
                    md_content.append(f"- **创建时间**: {core_props.created.strftime('%Y-%m-%d %H:%M:%S')}")
                if hasattr(core_props, 'modified') and core_props.modified:
                    md_content.append(f"- **修改时间**: {core_props.modified.strftime('%Y-%m-%d %H:%M:%S')}")
            except:
                pass
            
            md_content.append("")
            
            # 生成目录
            md_content.append("## 目录")
            md_content.append("")
            
            # 预先扫描标题以生成目录
            toc_items = []
            for idx, slide in enumerate(prs.slides, 1):
                try:
                    if slide.shapes.title and slide.shapes.title.text.strip():
                        title = slide.shapes.title.text.strip()
                        # 清理标题，移除换行符
                        title = ' '.join(title.split())
                        toc_items.append(f"{idx}. [{title}](#幻灯片-{idx})")
                    else:
                        toc_items.append(f"{idx}. [幻灯片 {idx}](#幻灯片-{idx})")
                except:
                    toc_items.append(f"{idx}. [幻灯片 {idx}](#幻灯片-{idx})")
            
            md_content.extend(toc_items)
            md_content.append("")
            
            # 转换每个幻灯片
            print(f"📊 共 {len(prs.slides)} 张幻灯片")
            
            for idx, slide in enumerate(prs.slides, 1):
                self.log(f"处理幻灯片 {idx}/{len(prs.slides)}")
                
                try:
                    slide_md = self.convert_slide_to_markdown(slide, idx, output_dir)
                    md_content.append(slide_md)
                except Exception as e:
                    self.log(f"转换幻灯片 {idx} 时出错: {e}")
                    md_content.append(f"\n---\n\n## 幻灯片 {idx}\n\n[转换失败]\n")
                
                # 显示进度
                if not self.verbose:
                    print(f"  进度: {idx}/{len(prs.slides)} ({idx*100//len(prs.slides)}%)", end='\r')
            
            if not self.verbose:
                print()  # 换行
            
            # 添加页脚
            md_content.append("\n---")
            md_content.append("")
            md_content.append(f"*此文档由 PPTX 转 Markdown 工具自动生成*")
            
            # 写入文件
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_content))
            
            # 输出统计信息
            print(f"\n✅ 转换成功！")
            print(f"📄 输出文件: {output_path}")
            print(f"📊 幻灯片数: {len(prs.slides)}")
            if self.extract_images and self.image_counter > 0:
                print(f"🖼️  提取图片: {self.image_counter}")
            if self.table_counter > 0:
                print(f"📋 提取表格: {self.table_counter}")
            
            return True
            
        except Exception as e:
            print(f"\n❌ 转换失败: {str(e)}")
            import traceback
            if self.verbose:
                traceback.print_exc()
            return False

def main():
    """主函数"""
    parser = argparse.ArgumentParser(
        description='将 PowerPoint (PPTX) 文件转换为 Markdown 格式（修复版）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
功能特点:
  • 提取文本、表格、图片
  • 保留基本格式（粗体、斜体、列表）
  • 生成文档目录
  • 提取幻灯片备注
  • 支持批量转换
  • 兼容不同版本的 python-pptx
  • 输出到与PPTX文件同名的文件夹

示例:
  %(prog)s presentation.pptx                    # 基本转换
  %(prog)s presentation.pptx -o output.md       # 指定输出文件
  %(prog)s *.pptx                               # 批量转换
  %(prog)s presentation.pptx --no-images        # 不提取图片
  %(prog)s presentation.pptx -v                 # 显示详细信息
        """
    )
    
    parser.add_argument('input', nargs='+', help='输入的 PPTX 文件（支持通配符）')
    parser.add_argument('-o', '--output', help='输出的 Markdown 文件路径')
    parser.add_argument('--no-images', action='store_true', help='不提取图片')
    parser.add_argument('-v', '--verbose', action='store_true', help='显示详细处理信息')
    
    args = parser.parse_args()
    
    # 创建转换器
    converter = PPTXToMarkdownConverter(
        extract_images=not args.no_images,
        verbose=args.verbose
    )
    
    # 处理输入文件
    input_files = []
    for pattern in args.input:
        if '*' in pattern or '?' in pattern:
            # 处理通配符
            from glob import glob
            files = glob(pattern)
            input_files.extend([f for f in files if f.lower().endswith('.pptx')])
        else:
            # 单个文件
            if Path(pattern).exists() and pattern.lower().endswith('.pptx'):
                input_files.append(pattern)
            else:
                print(f"⚠️  跳过无效文件: {pattern}")
    
    if not input_files:
        print("❌ 错误：没有找到有效的 PPTX 文件")
        sys.exit(1)
    
    # 转换文件
    success_count = 0
    
    if len(input_files) == 1 and args.output:
        # 单文件转换，指定了输出路径
        if converter.convert_presentation(input_files[0], args.output):
            success_count += 1
    else:
        # 批量转换或单文件转换（自动命名）
        for pptx_file in input_files:
            # 不再需要指定输出文件名，convert_presentation 会自动处理
            if converter.convert_presentation(pptx_file):
                success_count += 1
    
    # 总结
    if len(input_files) > 1:
        print(f"\n📊 总计: 成功转换 {success_count}/{len(input_files)} 个文件")

if __name__ == "__main__":
    main()

