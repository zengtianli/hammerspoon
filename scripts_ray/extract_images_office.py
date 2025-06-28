#!/usr/bin/env python3
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import subprocess
from pathlib import Path

def convert_wmf_files(directory):
    """调用 wmf2png.sh 脚本转换目录中的 WMF 文件"""
    # 检查目录中是否有 WMF 文件
    wmf_files = list(Path(directory).glob("*.wmf"))
    
    if wmf_files:
        print(f"  发现 {len(wmf_files)} 个 WMF 文件，尝试转换...")
        
        # convert_wmf_to_png.sh 脚本的路径
        script_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "convert_wmf_to_png.sh")
        
        try:
            # 调用 wmf2png.sh 脚本
            subprocess.run(["bash", script_path], cwd=directory, check=True)
            print("  WMF 文件转换完成")
        except subprocess.CalledProcessError as e:
            print(f"  WMF 转换出错: {e}")
        except Exception as e:
            print(f"  WMF 转换异常: {e}")

def extract_images_from_docx(docx_path, output_dir):
    """从 DOCX 文件中提取所有图片"""
    try:
        doc = Document(docx_path)
        
        # 创建输出目录
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # 获取文件名（不含扩展名）
        base_name = os.path.splitext(os.path.basename(docx_path))[0]
        
        # 提取图片
        img_count = 0
        has_wmf = False
        
        for i, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.target_ref:
                img = rel.target_part.blob
                img_count += 1
                
                # 获取图片扩展名
                ext = rel.target_ref.split('.')[-1].lower()
                
                # 检查是否是 WMF 文件
                if ext == 'wmf':
                    has_wmf = True
                    img_name = f"{base_name}_img_{img_count}.{ext}"
                elif ext not in ['png', 'jpg', 'jpeg', 'gif', 'bmp']:
                    ext = 'png'  # 默认使用 png
                    img_name = f"{base_name}_img_{img_count}.{ext}"
                else:
                    img_name = f"{base_name}_img_{img_count}.{ext}"
                
                img_path = os.path.join(output_dir, img_name)
                
                with open(img_path, "wb") as f:
                    f.write(img)
                
                print(f"  已保存: {img_path}")
        
        if img_count == 0:
            print(f"  未找到图片")
            # 如果没有图片，删除空文件夹
            if os.path.exists(output_dir) and not os.listdir(output_dir):
                os.rmdir(output_dir)
        else:
            print(f"  共提取 {img_count} 张图片")
            
            # 如果有WMF文件，调用转换脚本
            if has_wmf:
                convert_wmf_files(output_dir)
                
    except Exception as e:
        print(f"  处理出错: {e}")

def extract_images_from_pptx(pptx_path, output_dir):
    """从 PPTX 文件中提取所有图片"""
    try:
        prs = Presentation(pptx_path)
        
        # 创建输出目录
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # 获取文件名（不含扩展名）
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        
        img_count = 0
        has_wmf = False
        
        # 遍历每个幻灯片
        for slide_num, slide in enumerate(prs.slides, 1):
            # 遍历幻灯片中的每个形状
            for shape in slide.shapes:
                # 检查是否是图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    
                    # 获取图片
                    image = shape.image
                    
                    # 获取图片扩展名
                    ext = image.ext.lower()
                    
                    # 检查是否是 WMF 文件
                    if ext == 'wmf':
                        has_wmf = True
                    
                    # 生成文件名（包含幻灯片编号）
                    img_name = f"{base_name}_slide{slide_num}_img{img_count}.{ext}"
                    img_path = os.path.join(output_dir, img_name)
                    
                    # 保存图片
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    
                    print(f"  已保存: {img_path}")
                
                # 检查是否是组合形状（可能包含图片）
                elif hasattr(shape, "shapes"):
                    for shape_in_group in shape.shapes:
                        if shape_in_group.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            img_count += 1
                            image = shape_in_group.image
                            ext = image.ext.lower()
                            
                            # 检查是否是 WMF 文件
                            if ext == 'wmf':
                                has_wmf = True
                                
                            img_name = f"{base_name}_slide{slide_num}_img{img_count}.{ext}"
                            img_path = os.path.join(output_dir, img_name)
                            
                            with open(img_path, "wb") as f:
                                f.write(image.blob)
                            
                            print(f"  已保存: {img_path}")
        
        if img_count == 0:
            print(f"  未找到图片")
            # 如果没有图片，删除空文件夹
            if os.path.exists(output_dir) and not os.listdir(output_dir):
                os.rmdir(output_dir)
        else:
            print(f"  共提取 {img_count} 张图片")
            
            # 如果有WMF文件，调用转换脚本
            if has_wmf:
                convert_wmf_files(output_dir)
                
    except Exception as e:
        print(f"  处理出错: {e}")

def main():
    """主函数：处理当前目录下所有的 DOCX 和 PPTX 文件"""
    current_dir = Path.cwd()
    docx_files = list(current_dir.glob("*.docx"))
    pptx_files = list(current_dir.glob("*.pptx"))
    
    all_files = docx_files + pptx_files
    
    if not all_files:
        print("当前目录未找到 DOCX 或 PPTX 文件")
        return
    
    print(f"找到 {len(docx_files)} 个 DOCX 文件和 {len(pptx_files)} 个 PPTX 文件\n")
    
    for file_path in all_files:
        # 跳过临时文件（以~$开头的）
        if file_path.name.startswith('~$'):
            continue
            
        print(f"处理: {file_path.name}")
        
        # 创建对应的输出目录名（文件名_img）
        output_dir = current_dir / f"{file_path.stem}_img"
        
        # 根据文件类型提取图片
        if file_path.suffix.lower() == '.docx':
            extract_images_from_docx(str(file_path), str(output_dir))
        elif file_path.suffix.lower() == '.pptx':
            extract_images_from_pptx(str(file_path), str(output_dir))
        
        print()

if __name__ == "__main__":
    main()
